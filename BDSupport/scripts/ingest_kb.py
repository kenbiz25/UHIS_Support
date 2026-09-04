"""Run inside Docker: python3 scripts/ingest_kb.py"""
import os, sys, pathlib, logging
sys.path.insert(0, ".")
logging.basicConfig(level=logging.WARNING)

from config.settings import settings
from core.knowledge.store_faiss import FaissStore
from adapters.llm.openai_client import get_openai

try:
    from docx import Document as DocxDocument
except: DocxDocument = None
try:
    from pypdf import PdfReader
except:
    try: from PyPDF2 import PdfReader
    except: PdfReader = None
try:
    from pptx import Presentation
except: Presentation = None

def extract_text(path):
    try:
        if path.lower().endswith(".docx") and DocxDocument:
            doc = DocxDocument(path)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip()).strip()
        elif path.lower().endswith(".pdf") and PdfReader:
            reader = PdfReader(path)
            return "\n".join(p.extract_text() or "" for p in reader.pages).strip()
        elif path.lower().endswith(".pptx") and Presentation:
            prs = Presentation(path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
            return "\n".join(texts)
        elif path.lower().endswith((".txt", ".md")):
            return open(path, encoding="utf-8", errors="ignore").read().strip()
    except Exception as e:
        print(f"  ⚠️  {path}: {e}")
    return ""

def chunk_text(text, size=500, overlap=50):
    words = text.split()
    chunks, i = [], 0
    while i < len(words):
        chunk = " ".join(words[i:i+size])
        if chunk.strip(): chunks.append(chunk)
        i += size - overlap
    return chunks

files = [os.path.join("kb", f) for f in os.listdir("kb")
         if f.lower().endswith((".docx", ".pdf", ".pptx", ".txt", ".md"))]
print(f"Found {len(files)} files\n")

all_chunks, all_texts = [], []
for fpath in files:
    text = extract_text(fpath)
    if not text:
        print(f"  ⚠️  Empty/skipped: {fpath}")
        continue
    doc_id = os.path.basename(fpath)
    for i, ch in enumerate(chunk_text(text)):
        all_texts.append(ch)
        all_chunks.append({"text": ch, "metadata": {
            "document_id": doc_id, "chunk_id": i,
            "title": doc_id, "source_path": fpath
        }})
    print(f"  ✅ {doc_id}")

print(f"\nEmbedding {len(all_chunks)} chunks...")
client = get_openai()
vectors = []
for i in range(0, len(all_texts), 50):
    resp = client.embeddings.create(model=settings.EMBED_MODEL, input=all_texts[i:i+50])
    vectors.extend([d.embedding for d in resp.data])
    print(f"  {min(i+50, len(all_texts))}/{len(all_texts)} embedded")

for f in ["documents.pkl", "index.faiss", "meta.pkl"]:
    p = pathlib.Path(settings.FAISS_INDEX_DIR) / f
    if p.exists(): p.unlink()

store = FaissStore(dim=len(vectors[0]), index_dir=settings.FAISS_INDEX_DIR)
for vec, chunk in zip(vectors, all_chunks):
    store.upsert(content=chunk["text"], metadata=chunk["metadata"], embedding=vec)
store.save()
print(f"\n✅ KB indexed: {store.size()} documents")
